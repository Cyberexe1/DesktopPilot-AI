"""
File Writer Controller — creates files of any type and opens them in the correct app.
Supports: txt, docx, pptx, html, css, js, py, jsx, tsx, json, md, and more.
"""

import logging
import os
import subprocess

log = logging.getLogger(__name__)

_USER = os.environ.get("USERNAME", "User")

ALLOWED_DIRS = [
    os.path.expanduser("~/Desktop"),
    os.path.expanduser("~/Documents"),
    os.path.expanduser("~/Downloads"),
    rf"C:\Users\{_USER}\Desktop",
    rf"C:\Users\{_USER}\Documents",
    rf"C:\Users\{_USER}\Downloads",
    rf"C:\Users\{_USER}\Projects",
    "D:/Projects",
    "C:/Projects",
]

# Map file extensions to the app that should open them
EXT_TO_APP = {
    # Text editors
    ".txt":  "notepad",
    ".log":  "notepad",
    ".ini":  "notepad",
    # Microsoft Office
    ".docx": "word",
    ".doc":  "word",
    ".pptx": "powerpoint",
    ".ppt":  "powerpoint",
    ".xlsx": "excel",
    ".xls":  "excel",
    # Code files → VS Code
    ".html": "vscode",
    ".css":  "vscode",
    ".js":   "vscode",
    ".jsx":  "vscode",
    ".ts":   "vscode",
    ".tsx":  "vscode",
    ".py":   "vscode",
    ".java": "vscode",
    ".cpp":  "vscode",
    ".c":    "vscode",
    ".json": "vscode",
    ".md":   "vscode",
    ".yaml": "vscode",
    ".yml":  "vscode",
    ".xml":  "vscode",
    ".sql":  "vscode",
    ".sh":   "vscode",
    ".bat":  "vscode",
    ".env":  "vscode",
    ".gitignore": "vscode",
    # PDF
    ".pdf":  "default",
    # Images
    ".png":  "default",
    ".jpg":  "default",
    ".jpeg": "default",
}


def create_file(filename: str, content: str = "", directory: str = "", slides: int = 0) -> str:
    """
    Create a new file with content and open it in the appropriate application.
    - .txt → opens in Notepad
    - .docx/.pptx/.xlsx → opens in Word/PowerPoint/Excel
    - .html/.css/.js/.py/.jsx etc → opens in VS Code

    `slides` (optional): for .pptx, the number of content slides the user asked for.
    """
    # Default to Desktop
    if not directory or directory.lower() == "desktop":
        directory = os.path.expanduser("~/Desktop")
    elif directory.lower() == "documents":
        directory = os.path.expanduser("~/Documents")
    elif directory.lower() == "downloads":
        directory = os.path.expanduser("~/Downloads")

    # Create directory if it doesn't exist (for project scaffolding)
    if not os.path.exists(directory):
        try:
            os.makedirs(directory, exist_ok=True)
        except Exception as e:
            return f"Cannot create directory: {e}"

    # Safety check
    real_dir = os.path.realpath(directory)
    allowed = any(
        real_dir.startswith(os.path.realpath(d))
        for d in ALLOWED_DIRS if os.path.exists(d)
    )
    if not allowed:
        return f"Cannot create files in: {directory} (not in allowed directories)"

    # Sanitize filename (allow dots, dashes, underscores)
    safe_name = "".join(c for c in filename if c.isalnum() or c in '._-/ ')
    if not safe_name:
        return "Invalid filename"

    filepath = os.path.join(directory, safe_name)

    # Create parent directories if needed (for nested paths like src/App.jsx)
    parent = os.path.dirname(filepath)
    if parent and not os.path.exists(parent):
        os.makedirs(parent, exist_ok=True)

    try:
        ext = os.path.splitext(safe_name)[1].lower()

        if ext == ".docx":
            _create_docx(filepath, content)
        elif ext == ".pptx":
            # For presentations, enrich content if too sparse
            enriched = _enrich_pptx_content(content, safe_name, requested_slides=slides)
            _create_pptx(filepath, enriched)
        elif ext == ".xlsx":
            _create_xlsx(filepath, content)
        else:
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(content)

        log.info(f"Created: {filepath}")
        # Immediately add the new file to the index so it shows up in searches
        # without waiting for the next full rescan or file-watcher event.
        try:
            from indexer.file_indexer import insert_file as _idx_insert
            from datetime import datetime
            _idx_insert(safe_name, filepath, datetime.now().isoformat())
        except Exception:
            pass
    except Exception as e:
        return f"Failed to create file: {e}"

    # Open in the correct application
    result = _open_in_correct_app(filepath)
    return f"Created and opened: {filepath}"


def create_project(project_name: str, framework: str, directory: str = "") -> str:
    """
    Create a project by opening a terminal and running the REAL CLI commands.
    Opens a visible CMD window so user can see the setup happening.
    Supports: vite, nextjs, nodejs, python, django, fastapi, html
    """
    if not directory:
        directory = os.path.expanduser("~/Desktop")

    # Normalize directory paths
    directory = directory.replace("/", "\\")
    if not os.path.exists(directory):
        os.makedirs(directory, exist_ok=True)

    project_dir = os.path.join(directory, project_name)

    if os.path.exists(project_dir):
        _open_in_vscode(project_dir)
        return f"Project already exists, opened in VS Code: {project_dir}"

    fw = framework.lower()

    # Get the terminal commands for this framework
    commands = _get_project_commands(project_name, fw, directory)

    if not commands:
        # Fallback: manual scaffold for unknown frameworks
        os.makedirs(project_dir, exist_ok=True)
        _write(project_dir, "README.md", f"# {project_name}\n\nProject created by DesktopPilot AI\n")
        _open_in_vscode(project_dir)
        return f"Project '{project_name}' created (basic scaffold): {project_dir}"

    # Run all commands in a visible terminal window
    _run_in_visible_terminal(commands, directory)

    return f"Setting up {framework} project '{project_name}' in terminal at {directory}. VS Code will open when ready."


def _get_project_commands(name: str, framework: str, directory: str) -> str:
    """
    Return the full command string to set up a project.
    These run in a visible CMD window so the user sees progress.
    """
    project_path = os.path.join(directory, name)

    if framework in ("vite", "react", "react-vite"):
        return (
            f'cd /d "{directory}" && '
            f'npm create vite@latest {name} -- --template react && '
            f'cd {name} && '
            f'npm install && '
            f'code . && '
            f'echo Project ready! You can run: npm run dev'
        )

    elif framework in ("nextjs", "next", "next.js"):
        return (
            f'cd /d "{directory}" && '
            f'npx create-next-app@latest {name} --js --no-tailwind --no-eslint --app --use-npm && '
            f'cd {name} && '
            f'code . && '
            f'echo Project ready! You can run: npm run dev'
        )

    elif framework in ("nodejs", "node", "express"):
        return (
            f'cd /d "{directory}" && '
            f'mkdir {name} && cd {name} && '
            f'npm init -y && '
            f'npm install express && '
            f'code . && '
            f'echo Project ready! You can run: node index.js'
        )

    elif framework in ("django",):
        return (
            f'cd /d "{directory}" && '
            f'mkdir {name} && cd {name} && '
            f'python -m venv venv && '
            f'venv\\Scripts\\pip install django && '
            f'venv\\Scripts\\django-admin startproject {name} . && '
            f'code . && '
            f'echo Project ready! Activate venv then run: python manage.py runserver'
        )

    elif framework in ("fastapi", "fast-api"):
        return (
            f'cd /d "{directory}" && '
            f'mkdir {name} && cd {name} && '
            f'python -m venv venv && '
            f'venv\\Scripts\\pip install fastapi uvicorn && '
            f'code . && '
            f'echo Project ready! Activate venv then run: uvicorn main:app --reload'
        )

    elif framework in ("python", "flask"):
        return (
            f'cd /d "{directory}" && '
            f'mkdir {name} && cd {name} && '
            f'python -m venv venv && '
            f'venv\\Scripts\\pip install flask && '
            f'code . && '
            f'echo Project ready! Activate venv then run: python app.py'
        )

    elif framework in ("html", "static", "website"):
        # Just create folder structure and open in VS Code
        os.makedirs(project_path, exist_ok=True)
        _scaffold_html(project_path, name)
        _open_in_vscode(project_path)
        return ""  # No terminal needed

    elif framework in ("angular",):
        return (
            f'cd /d "{directory}" && '
            f'npx @angular/cli new {name} --skip-git --skip-tests && '
            f'cd {name} && '
            f'code . && '
            f'echo Project ready! You can run: ng serve'
        )

    elif framework in ("vue", "vue.js", "vuejs"):
        return (
            f'cd /d "{directory}" && '
            f'npm create vue@latest {name} -- --default && '
            f'cd {name} && '
            f'npm install && '
            f'code . && '
            f'echo Project ready! You can run: npm run dev'
        )

    elif framework in ("svelte", "sveltekit"):
        return (
            f'cd /d "{directory}" && '
            f'npm create svelte@latest {name} && '
            f'cd {name} && '
            f'npm install && '
            f'code . && '
            f'echo Project ready! You can run: npm run dev'
        )

    else:
        # Generic: create folder + npm init
        return (
            f'cd /d "{directory}" && '
            f'mkdir {name} && cd {name} && '
            f'npm init -y && '
            f'code . && '
            f'echo Project created!'
        )


def _run_in_visible_terminal(commands: str, cwd: str = "") -> None:
    """Open a CMD window and run commands visibly so the user can watch."""
    import tempfile, uuid
    try:
        # Write commands to a temp .bat file to avoid quoting issues with
        # long && chains passed through start/shell.
        bat = os.path.join(tempfile.gettempdir(), f"dp_{uuid.uuid4().hex[:8]}.bat")
        with open(bat, "w") as f:
            f.write("@echo off\n")
            f.write(commands + "\n")
            f.write("del /f /q \"%~f0\"\n")   # self-delete after running
        subprocess.Popen(
            f'start cmd /k "{bat}"',
            shell=True,
            cwd=cwd if cwd and os.path.exists(cwd) else None,
        )
        log.info(f"Opened terminal with: {commands[:100]}...")
    except Exception as e:
        log.error(f"Failed to open terminal: {e}")


def write_to_file(filepath: str, content: str, mode: str = "w") -> str:
    """Write content to an existing file."""
    if not os.path.exists(filepath):
        return f"File not found: {filepath}"
    try:
        with open(filepath, mode, encoding='utf-8') as f:
            f.write(content)
        return f"Written to: {filepath}"
    except Exception as e:
        return f"Failed to write: {e}"


# ── Internal helpers ──────────────────────────────────────────────────────────

def _open_in_correct_app(filepath: str) -> str:
    """Open a file in the correct application based on extension."""
    ext = os.path.splitext(filepath)[1].lower()
    app = EXT_TO_APP.get(ext, "default")

    if app == "vscode":
        _open_in_vscode(filepath)
    elif app == "notepad":
        subprocess.Popen(["notepad.exe", filepath], shell=False)
    elif app == "word":
        os.startfile(filepath)
    elif app == "powerpoint":
        os.startfile(filepath)
    elif app == "excel":
        os.startfile(filepath)
    else:
        os.startfile(filepath)

    return f"Opened: {filepath}"


def _open_in_vscode(path: str):
    """Open a file or folder in VS Code."""
    try:
        subprocess.Popen(["code", path], shell=False)
    except FileNotFoundError:
        vscode = rf"C:\Users\{_USER}\AppData\Local\Programs\Microsoft VS Code\Code.exe"
        try:
            subprocess.Popen([vscode, path], shell=False)
        except Exception:
            os.startfile(path)


def _write(directory: str, filename: str, content: str):
    """Helper to write a file inside a directory."""
    filepath = os.path.join(directory, filename)
    parent = os.path.dirname(filepath)
    if parent and not os.path.exists(parent):
        os.makedirs(parent, exist_ok=True)
    with open(filepath, 'w', encoding='utf-8') as f:
        f.write(content)


# ── Project scaffolders ───────────────────────────────────────────────────────

# ── Project scaffolders ───────────────────────────────────────────────────────

def _create_docx(filepath: str, content: str):
    """Create a professionally formatted .docx Word document."""
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # ── Page margins ──
    for section in doc.sections:
        section.top_margin    = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

    # ── Default body style ──
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    style.paragraph_format.space_after = Pt(6)

    lines = content.split('\n')

    def _is_heading(line: str) -> bool:
        """Detect heading lines — short, not starting with bullet, not lowercase start."""
        line = line.strip()
        if not line: return False
        if line.startswith(('•', '-', '*', '·', '\t', '   ')): return False
        if len(line) > 80: return False
        return True

    def _set_heading_style(para, level: int):
        """Apply styled heading formatting."""
        para.clear()
        run = para.add_run(para.text if hasattr(para, '_element') else '')
        colors = {1: RGBColor(0x1F, 0x39, 0x64), 2: RGBColor(0x2E, 0x74, 0xB5), 3: RGBColor(0x5B, 0x9B, 0xD5)}
        sizes  = {1: Pt(20), 2: Pt(16), 3: Pt(13)}
        para.style = f'Heading {level}'
        para.paragraph_format.space_before = Pt(14 if level == 1 else 10)
        para.paragraph_format.space_after  = Pt(4)

    first_line = True
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # ── Title (first non-empty line) ──
        if first_line and _is_heading(stripped):
            first_line = False
            para = doc.add_heading(stripped, level=0)
            para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = para.runs[0] if para.runs else para.add_run(stripped)
            run.font.color.rgb = RGBColor(0x1F, 0x39, 0x64)
            run.font.size = Pt(22)
            run.bold = True
            # Add a decorative line below title
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(0)
            p.paragraph_format.space_after  = Pt(12)
            run = p.add_run('─' * 60)
            run.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
            run.font.size = Pt(10)
            i += 1
            continue

        first_line = False

        # ── Section heading ──
        if _is_heading(stripped) and len(stripped) < 60:
            para = doc.add_heading(stripped, level=2)
            para.paragraph_format.space_before = Pt(12)
            para.paragraph_format.space_after  = Pt(4)
            if para.runs:
                para.runs[0].font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
            i += 1
            continue

        # ── Bullet point ──
        if stripped.startswith(('•', '-', '*', '·')):
            bullet_text = stripped.lstrip('•-*· ').strip()
            para = doc.add_paragraph(style='List Bullet')
            para.paragraph_format.left_indent  = Inches(0.3)
            para.paragraph_format.space_after  = Pt(3)
            run = para.add_run(bullet_text)
            run.font.name = 'Calibri'
            run.font.size = Pt(11)
            i += 1
            continue

        # ── "Dear Sir/Madam" → email greeting ──
        if stripped.lower().startswith('dear'):
            para = doc.add_paragraph()
            para.paragraph_format.space_after = Pt(12)
            run = para.add_run(stripped)
            run.font.name = 'Calibri'
            run.font.size = Pt(11)
            i += 1
            continue

        # ── "Best regards" → closing ──
        if any(stripped.lower().startswith(c) for c in ['best regards', 'sincerely', 'thank you', 'yours']):
            para = doc.add_paragraph()
            para.paragraph_format.space_before = Pt(12)
            run = para.add_run(stripped)
            run.font.name  = 'Calibri'
            run.font.size  = Pt(11)
            run.font.bold  = True
            i += 1
            continue

        # ── Regular paragraph ──
        para = doc.add_paragraph()
        para.paragraph_format.space_after = Pt(6)
        run = para.add_run(stripped)
        run.font.name = 'Calibri'
        run.font.size = Pt(11)
        i += 1

    doc.save(filepath)


def _enrich_pptx_content(content: str, filename: str, requested_slides: int = 0) -> str:
    """
    If AI-generated content is too sparse, call Bedrock AGAIN
    to generate detailed presentation content (plain text, not JSON).

    `requested_slides`: the number of content slides the user asked for.
    When set, it overrides the heuristic and content is regenerated to match.
    """
    import re

    lines = [l.strip() for l in content.split('\n') if l.strip()]

    # If content is already rich AND the user didn't ask for a specific count,
    # use it as-is. If a count was requested, always (re)generate to match it.
    bullet_count = sum(1 for l in lines if l.startswith('•'))
    if bullet_count >= 15 and not requested_slides:
        return content

    # Extract topic from filename or content
    topic = filename.replace('.pptx', '').replace('.', ' ').strip()
    if lines:
        topic = lines[0]  # First line is usually the topic

    # Check for flowchart request
    full_lower = content.lower()
    has_flowchart = any(kw in full_lower for kw in ['flowchart', 'flow chart', 'diagram', 'process flow'])

    # ── Determine number of slides ──
    # 1) explicit param from the planner (most reliable)
    # 2) a number parsed from the content text ("10 slides")
    # 3) heuristic from the number of heading-like lines
    if requested_slides and requested_slides > 0:
        num_slides = requested_slides
    else:
        m = re.search(r'(\d+)\s*slides?', full_lower)
        if m:
            num_slides = int(m.group(1))
        else:
            num_slides = max(4, len([l for l in lines if len(l) < 50 and not l.startswith('•')]))
    # Sane bounds: at least 1 content slide, cap to keep generation reasonable.
    num_slides = max(1, min(num_slides, 30))

    # Extra instructions from original content
    extra = ""
    if has_flowchart:
        extra += "Include a section listing process steps for a flowchart. "
    if lines[1:]:
        extra += f"Key points to cover: {', '.join(lines[1:5])}"

    # ── SECOND BEDROCK CALL — generate rich content ──
    from ai.content_generator import generate_content

    log.info(f"Calling Bedrock again for presentation content: '{topic}'")
    rich_content = generate_content(
        topic=topic,
        content_type="presentation",
        num_slides=num_slides,
        extra_instructions=extra,
    )

    # Append flowchart keyword if requested
    if has_flowchart and "flowchart" not in rich_content.lower():
        # Extract step-like lines for the flowchart
        flowchart_lines = [l for l in lines if len(l) < 40 and not l.startswith('•')]
        if len(flowchart_lines) > 2:
            rich_content += "\nflowchart\n" + "\n".join(flowchart_lines[1:])
        else:
            rich_content += "\nflowchart\nData Collection\nProcessing\nAnalysis\nAction\nMonitoring"

    return rich_content


def _create_pptx(filepath: str, content: str):
    """Create a professional dark-themed .pptx presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── Color palette (modern dark theme) ──
    DARK_BG    = RGBColor(0x0F, 0x29, 0x44)   # Deep navy
    ACCENT1    = RGBColor(0x00, 0xB0, 0xF0)   # Bright cyan
    ACCENT2    = RGBColor(0xFF, 0xC0, 0x00)   # Gold
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xCC, 0xDD, 0xEE)
    SLIDE_W    = Inches(13.33)
    SLIDE_H    = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    lines = [l.strip() for l in content.split('\n') if l.strip()]
    if not lines:
        lines = ["Presentation"]

    full_text = content.lower()
    has_flowchart = any(kw in full_text for kw in ['flowchart', 'flow chart', 'diagram', 'process flow'])

    content_lines = [l for l in lines if l.lower() not in ['flowchart', 'flow chart', 'diagram', 'process flow', 'steps:']]

    # ── Helper: dark background for any slide ──
    def _dark_bg(slide):
        from pptx.util import Inches
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    # ── Helper: add text box ──
    def _textbox(slide, text, left, top, width, height,
                 size=Pt(18), bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        return txBox

    # ── Helper: accent bar ──
    def _accent_bar(slide, color=ACCENT1, width_frac=0.15):
        from pptx.enum.shapes import MSO_SHAPE
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.33 * width_frac), SLIDE_H
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        return bar

    # ── Helper: horizontal rule ──
    def _hrule(slide, top, color=ACCENT1):
        from pptx.enum.shapes import MSO_SHAPE
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.2), top, Inches(10.9), Pt(3)
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = color
        rule.line.fill.background()

    # ── Parse slides data ──
    # ── Parse slides data (markdown-aware) ──
    def _strip_md(s: str) -> str:
        s = s.strip().lstrip('#').strip()
        return s.replace('**', '').replace('__', '').strip()

    def _is_separator(s: str) -> bool:
        s = s.strip()
        return len(s) >= 3 and set(s) <= set('-=*_~')

    title_text    = _strip_md(content_lines[0]) if content_lines else "Presentation"
    subtitle_text = _strip_md(content_lines[1]) if len(content_lines) > 1 else "Created by DesktopPilot AI"
    remaining     = content_lines[2:] if len(content_lines) > 2 else []

    slides_data = []
    current     = {"title": "", "bullets": []}
    for raw in remaining:
        line = raw.strip()
        if not line or _is_separator(line):
            continue

        md_heading = line.startswith('#')
        bullet     = line.startswith(('•', '·')) or (line.startswith(('-', '*', '+')) and not md_heading)
        text       = _strip_md(line) if md_heading else line.lstrip('•-*·+ ').strip()

        if bullet:
            if text:
                current["bullets"].append(text)
            continue

        # Heading: a markdown heading, or a short non-bullet line.
        is_head = md_heading or (len(line) < 80 and bool(text) and text[:1].isalnum())
        if is_head:
            if current["title"] or current["bullets"]:
                slides_data.append(current)
            current = {"title": text, "bullets": []}
        elif text:
            current["bullets"].append(text)

    if current["title"] or current["bullets"]:
        slides_data.append(current)

    if not slides_data and remaining:
        for i in range(0, len(remaining), 5):
            chunk = [_strip_md(c) for c in remaining[i:i+5]]
            slides_data.append({
                "title":   chunk[0],
                "bullets": chunk[1:] or ["Details for this section"]
            })

    # ════════════════════════════════════════════════════
    # SLIDE 1 — Title slide
    # ════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _dark_bg(slide)

    # Left accent bar
    _accent_bar(slide, ACCENT1, 0.04)

    # Large title
    _textbox(slide, title_text,
             Inches(1.0), Inches(2.0), Inches(11.0), Inches(2.0),
             size=Pt(40), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Horizontal rule
    _hrule(slide, Inches(4.2), ACCENT1)

    # Subtitle
    _textbox(slide, subtitle_text,
             Inches(1.0), Inches(4.5), Inches(11.0), Inches(1.0),
             size=Pt(20), color=ACCENT1, align=PP_ALIGN.CENTER, italic=True)

    # Footer
    _textbox(slide, "DesktopPilot AI",
             Inches(0.5), Inches(6.8), Inches(12.0), Inches(0.5),
             size=Pt(10), color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════
    # CONTENT SLIDES
    # ════════════════════════════════════════════════════
    for idx, sd in enumerate(slides_data[:30]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _dark_bg(slide)

        # Thin left accent bar
        _accent_bar(slide, ACCENT1 if idx % 2 == 0 else ACCENT2, 0.012)

        # Slide number badge
        _textbox(slide, f"{idx + 1:02d}",
                 Inches(11.8), Inches(0.15), Inches(1.2), Inches(0.5),
                 size=Pt(11), bold=True, color=ACCENT1, align=PP_ALIGN.RIGHT)

        # Section title
        _textbox(slide, sd["title"] or "Content",
                 Inches(0.5), Inches(0.25), Inches(11.5), Inches(0.9),
                 size=Pt(28), bold=True, color=WHITE)

        # Accent rule under title
        _hrule(slide, Inches(1.25), ACCENT1 if idx % 2 == 0 else ACCENT2)

        # Bullets
        bullets = sd["bullets"][:7]  # Max 7 bullets
        if bullets:
            bullet_top = Inches(1.5)
            spacing    = (SLIDE_H - bullet_top - Inches(0.5)) / max(len(bullets), 1)
            spacing    = min(spacing, Inches(0.75))

            for j, bullet in enumerate(bullets):
                y = bullet_top + (j * spacing)

                # Bullet dot
                from pptx.enum.shapes import MSO_SHAPE
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(0.45), y + Inches(0.15), Inches(0.12), Inches(0.12)
                )
                dot.fill.solid()
                dot.fill.fore_color.rgb = ACCENT1 if idx % 2 == 0 else ACCENT2
                dot.line.fill.background()

                # Bullet text
                _textbox(slide, bullet,
                         Inches(0.7), y, Inches(12.0), spacing,
                         size=Pt(16), color=LIGHT_GRAY)

        # Slide footer
        _textbox(slide, title_text,
                 Inches(0.5), Inches(7.15), Inches(12.0), Inches(0.3),
                 size=Pt(9), color=RGBColor(0x55, 0x77, 0x99), align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════
    # FLOWCHART SLIDE
    # ════════════════════════════════════════════════════
    if has_flowchart:
        steps = [sd["title"] for sd in slides_data if sd["title"]] or \
                ["Start", "Process", "Analysis", "Output", "End"]
        _add_flowchart_slide(prs, steps)

    # ════════════════════════════════════════════════════
    # THANK YOU SLIDE
    # ════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    _accent_bar(slide, ACCENT2, 0.04)

    _textbox(slide, "Thank You",
             Inches(1.0), Inches(2.5), Inches(11.0), Inches(1.5),
             size=Pt(48), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _hrule(slide, Inches(4.2), ACCENT2)

    _textbox(slide, "Questions & Discussion",
             Inches(1.0), Inches(4.5), Inches(11.0), Inches(0.8),
             size=Pt(20), color=ACCENT1, align=PP_ALIGN.CENTER, italic=True)

    prs.save(filepath)

    prs.save(filepath)


def _add_flowchart_slide(prs, steps: list):
    """Add a slide with a vertical flowchart that FITS within slide boundaries."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "Process Flow"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True

    # Layout constants — designed to fit within standard 10x7.5 inch slide
    max_steps = min(len(steps), 7)  # Max 7 steps to fit on one slide
    box_width = Inches(3.0)
    box_height = Inches(0.55)
    arrow_height = Inches(0.25)

    # Calculate total height needed and center vertically
    total_height = (max_steps * box_height) + ((max_steps - 1) * arrow_height)
    available_height = Inches(6.5)  # Slide height minus title and margins
    start_y = Inches(0.9) + (available_height - total_height) / 2
    start_x = (Inches(10) - box_width) / 2  # Center horizontally

    # Clamp start_y to minimum
    if start_y < Inches(0.8):
        start_y = Inches(0.8)

    # Colors for boxes
    colors = [
        RGBColor(0x4F, 0x8E, 0xF7),  # Blue
        RGBColor(0x22, 0xC5, 0x5E),  # Green
        RGBColor(0x7C, 0x5C, 0xFC),  # Purple
        RGBColor(0xF5, 0x9E, 0x0B),  # Orange
        RGBColor(0xEF, 0x44, 0x44),  # Red
        RGBColor(0x06, 0xB6, 0xD4),  # Cyan
        RGBColor(0xEC, 0x48, 0x99),  # Pink
    ]

    for i in range(max_steps):
        step_text = steps[i] if i < len(steps) else f"Step {i+1}"
        y = start_y + (i * (box_height + arrow_height))

        # Ensure we don't go off the slide (7.5 inches total height)
        if y + box_height > Inches(7.2):
            break

        # Draw rounded rectangle box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            start_x, y, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.color.rgb = colors[i % len(colors)]

        # Text inside box
        tf = shape.text_frame
        tf.word_wrap = True
        tf.text = step_text
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Arrow to next box (except last)
        if i < max_steps - 1:
            arrow_y = y + box_height
            arrow_x = start_x + (box_width / 2) - Inches(0.12)
            # Only add arrow if it fits
            if arrow_y + arrow_height < Inches(7.2):
                connector = slide.shapes.add_shape(
                    MSO_SHAPE.DOWN_ARROW,
                    arrow_x, arrow_y, Inches(0.24), arrow_height
                )
                connector.fill.solid()
                connector.fill.fore_color.rgb = RGBColor(0x64, 0x74, 0x8B)
                connector.line.fill.background()


def _create_xlsx(filepath: str, content: str):
    """Create a proper .xlsx Excel spreadsheet."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"

    # Split content into rows (by newline) and columns (by comma or tab)
    for row_idx, line in enumerate(content.split('\n'), start=1):
        if not line.strip():
            continue
        # Try comma-separated first, then tab, then single cell
        if ',' in line:
            cells = line.split(',')
        elif '\t' in line:
            cells = line.split('\t')
        else:
            cells = [line]

        for col_idx, cell in enumerate(cells, start=1):
            ws.cell(row=row_idx, column=col_idx, value=cell.strip())

    wb.save(filepath)


def _scaffold_vite(d: str, name: str):
    _write(d, "package.json", f'''{{"name": "{name}", "version": "1.0.0", "scripts": {{"dev": "vite", "build": "vite build"}}, "dependencies": {{"react": "^18.3.1", "react-dom": "^18.3.1"}}, "devDependencies": {{"vite": "^5.3.1", "@vitejs/plugin-react": "^4.3.1"}}}}''')
    _write(d, "index.html", f'''<!DOCTYPE html>\n<html lang="en">\n<head>\n  <meta charset="UTF-8" />\n  <title>{name}</title>\n</head>\n<body>\n  <div id="root"></div>\n  <script type="module" src="/src/main.jsx"></script>\n</body>\n</html>''')
    _write(d, "vite.config.js", "import { defineConfig } from 'vite'\nimport react from '@vitejs/plugin-react'\n\nexport default defineConfig({ plugins: [react()] })\n")
    _write(d, "src/main.jsx", "import React from 'react'\nimport ReactDOM from 'react-dom/client'\nimport App from './App'\n\nReactDOM.createRoot(document.getElementById('root')).render(<App />)\n")
    _write(d, "src/App.jsx", f"export default function App() {{\n  return <h1>{name}</h1>\n}}\n")


def _scaffold_nextjs(d: str, name: str):
    _write(d, "package.json", f'''{{"name": "{name}", "version": "1.0.0", "scripts": {{"dev": "next dev", "build": "next build", "start": "next start"}}, "dependencies": {{"next": "^14.0.0", "react": "^18.3.1", "react-dom": "^18.3.1"}}}}''')
    _write(d, "app/page.jsx", f"export default function Home() {{\n  return <h1>Welcome to {name}</h1>\n}}\n")
    _write(d, "app/layout.jsx", f"export default function RootLayout({{ children }}) {{\n  return <html><body>{{children}}</body></html>\n}}\n")


def _scaffold_nodejs(d: str, name: str):
    _write(d, "package.json", f'''{{"name": "{name}", "version": "1.0.0", "main": "index.js", "scripts": {{"start": "node index.js", "dev": "nodemon index.js"}}, "dependencies": {{"express": "^4.18.2"}}}}''')
    _write(d, "index.js", "const express = require('express')\nconst app = express()\n\napp.get('/', (req, res) => res.send('Hello World'))\n\napp.listen(3000, () => console.log('Server running on port 3000'))\n")


def _scaffold_python(d: str, name: str):
    _write(d, "requirements.txt", "fastapi==0.111.0\nuvicorn==0.30.1\n")
    _write(d, "main.py", "from fastapi import FastAPI\n\napp = FastAPI()\n\n@app.get('/')\ndef root():\n    return {'message': 'Hello World'}\n")
    _write(d, "README.md", f"# {name}\n\nRun: `uvicorn main:app --reload`\n")


def _scaffold_html(d: str, name: str):
    _write(d, "index.html", f"<!DOCTYPE html>\n<html>\n<head>\n  <title>{name}</title>\n  <link rel='stylesheet' href='style.css'>\n</head>\n<body>\n  <h1>{name}</h1>\n  <script src='script.js'></script>\n</body>\n</html>\n")
    _write(d, "style.css", "body { font-family: sans-serif; margin: 2rem; }\nh1 { color: #333; }\n")
    _write(d, "script.js", "console.log('Hello from " + name + "')\n")


